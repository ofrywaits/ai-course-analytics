# -*- coding: utf-8 -*-
"""
Generate a professional business presentation (.pptx) for the
AI Course Analytics Platform final project.

Usage:
    python generate_presentation.py
    → outputs/presentation.pptx
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

logger = logging.getLogger(__name__)

# ── Design system (matches DESIGN.md / Stitch prototype) ─────────────────────
BG       = RGBColor(0x0E, 0x14, 0x17)   # #0E1417  background
SURFACE  = RGBColor(0x1A, 0x21, 0x23)   # #1A2123  card surface
BORDER   = RGBColor(0x3C, 0x49, 0x4E)   # #3C494E  outline-variant
CYAN     = RGBColor(0x00, 0xD4, 0xFF)   # #00D4FF  primary accent
WHITE    = RGBColor(0xDD, 0xE3, 0xE7)   # #DDE3E7  on-surface
MUTED    = RGBColor(0xBB, 0xC9, 0xCF)   # #BBC9CF  on-surface-variant
AMBER    = RGBColor(0xFF, 0xB5, 0x28)   # tertiary
GREEN    = RGBColor(0x4A, 0xDE, 0x80)   # success green

# Slide dimensions: widescreen 16:9
W = Inches(13.33)
H = Inches(7.5)

SLIDE_W = 13.33   # inches
SLIDE_H = 7.5

OUT_PATH = Path("outputs/presentation.pptx")
PROTO_DIR = Path("docs/ux-prototype")


# ── Helpers ──────────────────────────────────────────────────────────────────

def rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a filled rectangle shape."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             font_name="Calibri"):
    """Add a text box."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or WHITE
    run.font.name = font_name
    return txBox


def add_label(slide, text, left, top, width, height,
              font_size=11, color=None):
    """Small uppercase label."""
    return add_text(slide, text.upper(), left, top, width, height,
                    font_size=font_size, bold=True,
                    color=color or MUTED, align=PP_ALIGN.LEFT)


def set_bg(slide):
    """Set slide background to dark color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def cyan_bar(slide, height=0.06):
    """Thin cyan accent bar at top of slide."""
    add_rect(slide, 0, 0, SLIDE_W, height, fill_color=CYAN)


def slide_number(slide, num):
    """Small slide number bottom-right."""
    add_text(slide, str(num), 12.6, 7.1, 0.5, 0.3,
             font_size=9, color=BORDER, align=PP_ALIGN.RIGHT)


def metric_card(slide, left, top, w, h, label, value, sub=None, value_color=None):
    """Metric card: surface bg + border + label + big value."""
    add_rect(slide, left, top, w, h, fill_color=SURFACE, border_color=BORDER, border_width=Pt(0.75))
    add_label(slide, label, left + 0.15, top + 0.12, w - 0.3, 0.3, font_size=9)
    add_text(slide, value, left + 0.1, top + 0.38, w - 0.2, 0.7,
             font_size=28, bold=True, color=value_color or CYAN, align=PP_ALIGN.CENTER)
    if sub:
        add_text(slide, sub, left + 0.1, top + 0.98, w - 0.2, 0.3,
                 font_size=10, color=MUTED, align=PP_ALIGN.CENTER)


def add_image(slide, path: Path, left, top, width, height):
    """Add image if file exists."""
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                                 Inches(width), Inches(height))


def bullet_block(slide, items: list[tuple[str, str]], left, top, width,
                 title=None, title_color=None):
    """Render a list of (emoji/bullet, text) items."""
    y = top
    if title:
        add_text(slide, title, left, y, width, 0.45,
                 font_size=13, bold=True, color=title_color or CYAN)
        y += 0.45
    for icon, line in items:
        add_text(slide, f"{icon}  {line}", left, y, width, 0.38,
                 font_size=12, color=WHITE)
        y += 0.38
    return y


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs: Presentation) -> None:
    """Slide 1 — Title."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(sl)
    # Left cyan accent strip
    add_rect(sl, 0, 0, 0.08, SLIDE_H, fill_color=CYAN)
    # Main title
    add_text(sl, "AI Course Analytics Platform",
             0.5, 1.4, 9, 1.1, font_size=40, bold=True, color=WHITE)
    # Cyan underline
    add_rect(sl, 0.5, 2.55, 5.5, 0.05, fill_color=CYAN)
    # Subtitle
    add_text(sl, "Automated Analysis of 200,000+ Udemy Courses\nUsing CrewAI Flow · scikit-learn · Streamlit",
             0.5, 2.7, 9, 0.9, font_size=18, color=MUTED)
    # Stats strip
    for i, (val, lbl) in enumerate([
        ("209K", "Courses Analyzed"),
        ("87.4%", "Model Accuracy"),
        ("6 AI Agents", "CrewAI Flow"),
        ("77 Tests", "79% Coverage"),
    ]):
        x = 0.5 + i * 3.1
        metric_card(sl, x, 3.8, 2.8, 1.35, lbl, val)
    # Bottom
    add_text(sl, "Final Project  ·  2026  ·  Ofry Waits",
             0.5, 6.8, 9, 0.45, font_size=11, color=BORDER)
    slide_number(sl, 1)


def slide_02_problem(prs: Presentation) -> None:
    """Slide 2 — The Problem."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "The Problem", 0.6, 0.25, 8, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    add_text(sl,
             '"I want to create an online course on Udemy — but where do I even start?"',
             0.6, 1.2, 11.8, 0.8, font_size=16, color=AMBER)

    # Pain points
    pains = [
        ("❓", "209,000+ courses on Udemy — impossible to analyse manually"),
        ("❓", "No data-driven way to choose the right category, price, or launch time"),
        ("❓", "Course creators make decisions based on gut feeling, not evidence"),
        ("❓", "Business question: What makes a course popular? Can we predict it?"),
    ]
    bullet_block(sl, pains, 0.6, 2.2, 11.5)

    # Maya persona card
    add_rect(sl, 0.6, 4.6, 12.1, 2.2, fill_color=SURFACE, border_color=CYAN, border_width=Pt(1.5))
    add_text(sl, "👤  Maya — Our User Persona", 0.85, 4.75, 8, 0.45,
             font_size=13, bold=True, color=CYAN)
    add_text(sl,
             "Independent educator, 5 years of experience. She wants to publish her first Udemy course "
             "on Python for Data Science — but has no idea whether to price it at $9.99 or $49.99, "
             "which subcategory to list under, or whether November or March is a better launch month. "
             "She needs a tool that tells her exactly what the data says.",
             0.85, 5.25, 11.6, 1.3, font_size=12, color=WHITE)
    slide_number(sl, 2)


def slide_03_solution(prs: Presentation) -> None:
    """Slide 3 — Our Solution."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Our Solution", 0.6, 0.25, 8, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    add_text(sl, "A fully automated AI pipeline that ingests raw data, generates insights, "
                 "trains a predictive model, and serves everything through an interactive dashboard.",
             0.6, 1.15, 12, 0.65, font_size=14, color=MUTED)

    # Architecture flow
    boxes = [
        (0.4,  "📄 Raw Data\n209,735\ncourses", BG),
        (2.8,  "🔵 Crew 1\nData Analyst\n3 AI Agents", SURFACE),
        (5.2,  "✅ Validation\nGate 1\nContract check", SURFACE),
        (7.6,  "🟣 Crew 2\nData Scientist\n3 AI Agents", SURFACE),
        (10.0, "🖥️ Streamlit\nDashboard\n5 Tabs", SURFACE),
    ]
    for x, lbl, bg in boxes:
        add_rect(sl, x, 2.1, 2.2, 1.6, fill_color=bg, border_color=CYAN, border_width=Pt(1))
        add_text(sl, lbl, x + 0.1, 2.2, 2.0, 1.4,
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrows
    for ax in [2.6, 5.0, 7.4, 9.8]:
        add_text(sl, "→", ax, 2.65, 0.3, 0.5, font_size=20, color=CYAN, align=PP_ALIGN.CENTER)

    # Outputs
    outputs = [
        ("Crew 1 outputs", ["clean_data.csv", "eda_report.html", "insights.md", "dataset_contract.json"], 0.6),
        ("Crew 2 outputs", ["features.csv", "model.pkl", "evaluation_report.md", "model_card.md"], 6.8),
    ]
    for title, items, x in outputs:
        add_text(sl, title, x, 4.05, 5.5, 0.35, font_size=11, bold=True, color=CYAN)
        for i, item in enumerate(items):
            col = x + (i % 2) * 2.7
            row = 4.45 + (i // 2) * 0.38
            add_text(sl, f"📄 {item}", col, row, 2.6, 0.35, font_size=10, color=WHITE)

    slide_number(sl, 3)


def slide_04_data(prs: Presentation) -> None:
    """Slide 4 — The Dataset."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "The Dataset", 0.6, 0.25, 8, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # Dataset metrics
    cards = [
        ("Raw Rows", "209,735", "Udemy courses"),
        ("After Cleaning", "206,885", "99% retention"),
        ("Features", "15+", "engineered features"),
        ("Time Range", "2011–2017", "6 years of data"),
    ]
    for i, (lbl, val, sub) in enumerate(cards):
        metric_card(sl, 0.4 + i * 3.1, 1.2, 2.8, 1.5, lbl, val, sub)

    # Cleaning pipeline
    add_text(sl, "7-Step Data Cleaning Pipeline", 0.6, 3.0, 9, 0.45,
             font_size=14, bold=True, color=CYAN)
    steps = [
        "1  Load raw CSV — validate schema and row count",
        "2  Remove duplicates — dropped 2,850 exact duplicates",
        "3  Fix data types — parse dates, cast numerics",
        "4  Handle missing values — subject/language filled with mode",
        "5  Remove outliers — top 0.5% subscribers trimmed",
        "6  Engineer features — title_length, publish_month, is_new",
        "7  Save clean_data.csv + dataset_contract.json",
    ]
    for i, s in enumerate(steps):
        col = 0.6 if i < 4 else 6.8
        row = 3.55 + (i % 4) * 0.72
        add_rect(sl, col, row, 5.8, 0.6, fill_color=SURFACE, border_color=BORDER)
        add_text(sl, s, col + 0.15, row + 0.1, 5.5, 0.42, font_size=11, color=WHITE)

    slide_number(sl, 4)


def slide_05_eda(prs: Presentation) -> None:
    """Slide 5 — EDA & Visual Insights."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Exploratory Data Analysis", 0.6, 0.25, 9, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # EDA screenshot left
    add_image(sl, PROTO_DIR / "03-eda-report.png", 0.4, 1.15, 6.0, 5.9)

    # Key findings right
    findings = [
        ("📈", "Development is #1 category\nwith 65K+ courses"),
        ("💰", "Free courses average 5,310\nsubscribers vs 1,779 paid"),
        ("📅", "November is best launch\nmonth for visibility"),
        ("⭐", "Average course rating: 3.74\n— quality matters"),
        ("🔗", "Reviews & comments are\nstrongest popularity signals"),
        ("📊", "6 charts embedded in\neda_report.html"),
    ]
    y = 1.15
    for icon, text in findings:
        add_rect(sl, 6.8, y, 6.1, 0.88, fill_color=SURFACE, border_color=BORDER)
        add_text(sl, f"{icon}  {text}", 7.0, y + 0.1, 5.7, 0.72, font_size=11, color=WHITE)
        y += 0.96

    slide_number(sl, 5)


def slide_06_insights(prs: Presentation) -> None:
    """Slide 6 — Business Insights."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Business Insights & Recommendations",
             0.6, 0.25, 12, 0.65, font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    recommendations = [
        ("🎯", "Focus on Development",
         "Highest demand category — 3× more courses than Business or Design. "
         "Sub-topics like Python, Web Dev, and Data Science dominate."),
        ("🆓", "Start Free, Then Monetize",
         "Free courses attract 3× more subscribers. Use a free intro course "
         "to build audience and upsell a paid advanced version."),
        ("📅", "Launch in November",
         "Holiday season drives a measurable spike in enrollments. "
         "Publish 2 weeks before to gain organic traction before peak."),
        ("⭐", "Quality > Quantity",
         "Courses rated 4.5+ get exponentially more reviews and subscribers. "
         "Invest in production quality: audio, structure, exercises."),
    ]
    for i, (icon, title, body) in enumerate(recommendations):
        x = 0.4 + (i % 2) * 6.4
        y = 1.3 + (i // 2) * 2.7
        add_rect(sl, x, y, 6.1, 2.45, fill_color=SURFACE, border_color=CYAN, border_width=Pt(1))
        add_text(sl, f"{icon}  {title}", x + 0.2, y + 0.15, 5.7, 0.45,
                 font_size=14, bold=True, color=CYAN)
        add_text(sl, body, x + 0.2, y + 0.65, 5.7, 1.65, font_size=12, color=WHITE)

    slide_number(sl, 6)


def slide_07_model(prs: Presentation) -> None:
    """Slide 7 — The Predictive Model."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Predictive Model — Course Popularity",
             0.6, 0.25, 10, 0.65, font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # Model comparison cards
    add_text(sl, "Model Comparison", 0.6, 1.15, 5, 0.4, font_size=14, bold=True, color=CYAN)
    models = [
        ("🌲 Random Forest", "87.4%", "87.4%", "Winner ✅", GREEN),
        ("📉 Logistic Regression", "86.6%", "86.5%", "Baseline", MUTED),
    ]
    for i, (name, acc, cv, badge, col) in enumerate(models):
        y = 1.65 + i * 1.55
        add_rect(sl, 0.5, y, 5.9, 1.35, fill_color=SURFACE, border_color=col, border_width=Pt(1.5))
        add_text(sl, name, 0.75, y + 0.1, 3.5, 0.4, font_size=13, bold=True, color=col)
        add_text(sl, f"Accuracy: {acc}   CV: {cv}   {badge}",
                 0.75, y + 0.6, 5.5, 0.4, font_size=12, color=WHITE)

    # Feature importances
    add_text(sl, "Top 5 Feature Importances", 0.6, 4.55, 5, 0.4,
             font_size=14, bold=True, color=CYAN)
    features = [
        ("num_reviews", 39.5),
        ("num_comments", 26.8),
        ("avg_rating", 6.9),
        ("price", 5.2),
        ("title_length", 4.1),
    ]
    for i, (feat, imp) in enumerate(features):
        y = 5.05 + i * 0.44
        add_text(sl, feat, 0.6, y, 2.5, 0.38, font_size=11, color=WHITE)
        bar_w = imp / 100 * 3.2
        add_rect(sl, 3.2, y + 0.08, bar_w, 0.22, fill_color=CYAN)
        add_text(sl, f"{imp}%", 3.2 + bar_w + 0.1, y, 0.6, 0.35, font_size=10, color=MUTED)

    # Model results screenshot right
    add_image(sl, PROTO_DIR / "05-model-results.png", 6.8, 1.15, 6.1, 5.9)

    slide_number(sl, 7)


def slide_08_dashboard(prs: Presentation) -> None:
    """Slide 8 — Streamlit Dashboard."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Streamlit Dashboard", 0.6, 0.25, 8, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # Overview screenshot (large, left)
    add_image(sl, PROTO_DIR / "01-overview.png", 0.4, 1.15, 7.5, 5.9)

    # Tab list right
    add_text(sl, "5 Interactive Tabs", 8.3, 1.15, 4.7, 0.45,
             font_size=14, bold=True, color=CYAN)
    tabs = [
        ("1", "Overview", "Health banner, KPI cards, run status"),
        ("2", "EDA Report", "6 embedded interactive charts"),
        ("3", "Model Results", "Accuracy, CV, confusion matrix"),
        ("4", "Course Predictor", "Live ML prediction for your course"),
        ("5", "Downloads", "All 8 output artifacts"),
    ]
    for i, (num, tab, desc) in enumerate(tabs):
        y = 1.7 + i * 1.1
        add_rect(sl, 8.2, y, 4.8, 0.98, fill_color=SURFACE, border_color=BORDER)
        add_text(sl, f"Tab {num} — {tab}", 8.4, y + 0.07, 4.4, 0.38,
                 font_size=12, bold=True, color=CYAN)
        add_text(sl, desc, 8.4, y + 0.52, 4.4, 0.35, font_size=10, color=MUTED)

    slide_number(sl, 8)


def slide_09_predictor(prs: Presentation) -> None:
    """Slide 9 — Course Predictor (Maya demo)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Course Popularity Predictor", 0.6, 0.25, 9, 0.65,
             font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # Left — predictor screenshot
    add_image(sl, PROTO_DIR / "02-course-predictor.png", 0.4, 1.15, 6.2, 5.9)

    # Right — Maya demo scenario
    add_rect(sl, 7.0, 1.15, 5.95, 5.9, fill_color=SURFACE, border_color=CYAN, border_width=Pt(1.5))
    add_text(sl, "👤  Maya's Demo Scenario", 7.2, 1.3, 5.5, 0.45,
             font_size=13, bold=True, color=CYAN)

    inputs = [
        ("Category", "Development"),
        ("Subcategory", "Data Science"),
        ("Price", "$0 (Free)"),
        ("Language", "English"),
        ("Title length", "52 characters"),
        ("Publish month", "November"),
    ]
    add_text(sl, "Input Parameters:", 7.2, 1.85, 5.5, 0.35,
             font_size=11, bold=True, color=MUTED)
    for i, (k, v) in enumerate(inputs):
        y = 2.25 + i * 0.48
        add_text(sl, k, 7.2, y, 2.2, 0.38, font_size=11, color=MUTED)
        add_text(sl, v, 9.5, y, 3.2, 0.38, font_size=11, bold=True, color=WHITE)

    # Result
    add_rect(sl, 7.2, 5.3, 5.5, 1.5, fill_color=BG, border_color=GREEN, border_width=Pt(2))
    add_text(sl, "🎯  Prediction Result", 7.4, 5.4, 5.0, 0.38,
             font_size=12, bold=True, color=GREEN)
    add_text(sl, "POPULAR  —  92% confidence", 7.4, 5.82, 5.0, 0.45,
             font_size=16, bold=True, color=WHITE)

    slide_number(sl, 9)


def slide_10_tech(prs: Presentation) -> None:
    """Slide 10 — Tech Stack & Architecture Decisions."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Tech Stack & Architecture Decisions",
             0.6, 0.25, 10, 0.65, font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    stack = [
        ("⚙️  AI Orchestration", "CrewAI Flow", "@start/@listen, sequential crews with validation gates"),
        ("🧠  LLM", "Groq llama-3.3-70b", "Free tier — Python computes, LLM only narrates (<300 words)"),
        ("📊  ML", "scikit-learn RF + LR", "Interpretable, no GPU, feature importances for dashboard"),
        ("🖥️  Dashboard", "Streamlit", "Zero front-end code, 5 tabs, built-in file downloads"),
        ("📁  MCP", "filesystem server", "Agents 3 & 6 read output files directly — grounded responses"),
        ("🚀  Deployment", "Railway", "Zero-config, auto-deploy on git push to main"),
    ]
    for i, (layer, tech, why) in enumerate(stack):
        x = 0.4 + (i % 2) * 6.4
        y = 1.3 + (i // 2) * 2.0
        add_rect(sl, x, y, 6.1, 1.8, fill_color=SURFACE, border_color=BORDER)
        add_text(sl, layer, x + 0.2, y + 0.1, 5.5, 0.38, font_size=10, color=MUTED)
        add_text(sl, tech, x + 0.2, y + 0.48, 5.5, 0.42, font_size=14, bold=True, color=CYAN)
        add_text(sl, why, x + 0.2, y + 0.95, 5.5, 0.7, font_size=11, color=WHITE)

    slide_number(sl, 10)


def slide_11_quality(prs: Presentation) -> None:
    """Slide 11 — Quality & Testing."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    cyan_bar(sl)
    add_text(sl, "Quality, Testing & GitHub Workflow",
             0.6, 0.25, 10, 0.65, font_size=28, bold=True, color=WHITE)
    add_rect(sl, 0.6, 0.95, 1.2, 0.04, fill_color=CYAN)

    # Test metrics
    test_cards = [
        ("Total Tests", "77", "all passing"),
        ("Coverage", "79%", "tools + validation"),
        ("Pull Requests", "3", "merged on GitHub"),
        ("Test Files", "7", "pytest modules"),
    ]
    for i, (lbl, val, sub) in enumerate(test_cards):
        metric_card(sl, 0.4 + i * 3.1, 1.2, 2.8, 1.5, lbl, val, sub)

    # Test files
    add_text(sl, "Test Coverage by Module", 0.6, 2.95, 6, 0.4,
             font_size=13, bold=True, color=CYAN)
    test_files = [
        ("test_data_tools.py", "Data pipeline: cleaning, outliers, features"),
        ("test_model_tools.py", "ML pipeline: training, evaluation, persistence"),
        ("test_validators.py", "Validation gates: Crew 1 & 2 contract checks"),
        ("test_monitoring.py", "Health checks, run metrics parsing"),
        ("test_mcp_tools.py", "MCP graceful fallback"),
        ("test_viz_tools.py", "All 6 chart functions, base64 output"),
        ("test_integration.py", "End-to-end: data → features → model"),
    ]
    for i, (f, desc) in enumerate(test_files):
        y = 3.45 + i * 0.54
        add_text(sl, f"✅  {f}", 0.6, y, 4.0, 0.42, font_size=11, bold=True, color=CYAN)
        add_text(sl, desc, 4.7, y, 8.0, 0.42, font_size=11, color=WHITE)

    # GitHub workflow
    add_text(sl, "GitHub Workflow", 0.6, 7.05, 12, 0.35,
             font_size=10, color=MUTED)

    slide_number(sl, 11)


def slide_12_summary(prs: Presentation) -> None:
    """Slide 12 — Summary & Next Steps."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    # Full-width cyan bar top
    add_rect(sl, 0, 0, SLIDE_W, 0.08, fill_color=CYAN)

    add_text(sl, "Summary & What's Next",
             0.6, 0.2, 9, 0.7, font_size=32, bold=True, color=WHITE)

    # Achievements
    add_text(sl, "✅  What We Built", 0.6, 1.1, 6, 0.42,
             font_size=14, bold=True, color=CYAN)
    achievements = [
        "Full CrewAI Flow pipeline — 2 Crews, 6 Agents, 2 validation gates",
        "8 output artifacts — data, EDA, ML model, model card",
        "87.4% accuracy Random Forest — beats Logistic Regression baseline",
        "Interactive Streamlit dashboard with live course popularity predictor",
        "77 tests, 79% coverage, MCP filesystem integration, Railway deployment",
    ]
    for i, a in enumerate(achievements):
        add_text(sl, f"•  {a}", 0.6, 1.62 + i * 0.5, 11.8, 0.42,
                 font_size=12, color=WHITE)

    # Lessons learned
    add_rect(sl, 0.5, 4.3, 5.8, 2.85, fill_color=SURFACE, border_color=BORDER)
    add_text(sl, "💡  Lessons Learned", 0.7, 4.45, 5.4, 0.4,
             font_size=13, bold=True, color=AMBER)
    lessons = [
        "LLM rate limits — Python-first design saved the project",
        "Validation gates prevent silent failures downstream",
        "50 tests caught 7 real bugs during audit phase",
        "MCP graceful fallback is essential for portability",
    ]
    for i, l in enumerate(lessons):
        add_text(sl, f"→  {l}", 0.7, 4.95 + i * 0.52, 5.4, 0.42,
                 font_size=11, color=WHITE)

    # v2 roadmap
    add_rect(sl, 6.9, 4.3, 5.9, 2.85, fill_color=SURFACE, border_color=BORDER)
    add_text(sl, "🚀  v2.0 Roadmap", 7.1, 4.45, 5.4, 0.4,
             font_size=13, bold=True, color=CYAN)
    roadmap = [
        ("High", "Real-time Udemy API integration"),
        ("High", "A/B price simulator"),
        ("Med", "Email report export for teams"),
        ("Med", "Category trend analysis (time series)"),
    ]
    for i, (pri, item) in enumerate(roadmap):
        col = CYAN if pri == "High" else AMBER
        add_text(sl, f"[{pri}]", 7.1, 5.0 + i * 0.52, 0.9, 0.38,
                 font_size=10, bold=True, color=col)
        add_text(sl, item, 8.1, 5.0 + i * 0.52, 3.5, 0.38, font_size=11, color=WHITE)

    # Footer
    add_rect(sl, 0, 7.15, SLIDE_W, 0.35, fill_color=SURFACE)
    add_text(sl, "github.com/ofrywaits/ai-course-analytics",
             0.5, 7.18, 7, 0.28, font_size=10, color=CYAN)
    add_text(sl, "Ofry Waits  ·  Final Project 2026",
             9.0, 7.18, 4, 0.28, font_size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    slide_number(sl, 12)


# ── Main ──────────────────────────────────────────────────────────────────────

def build_presentation() -> Path:
    """Build the full 12-slide presentation and save to outputs/."""
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    builders = [
        slide_01_title,
        slide_02_problem,
        slide_03_solution,
        slide_04_data,
        slide_05_eda,
        slide_06_insights,
        slide_07_model,
        slide_08_dashboard,
        slide_09_predictor,
        slide_10_tech,
        slide_11_quality,
        slide_12_summary,
    ]

    for i, build_fn in enumerate(builders, 1):
        logger.info("Building slide %d — %s", i, build_fn.__doc__.split("—")[0].strip())
        build_fn(prs)

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO, format="%(message)s")
    path = build_presentation()
    print(f"\n✅  Presentation saved → {path}")
    print(f"   Open with: open '{path}'")
